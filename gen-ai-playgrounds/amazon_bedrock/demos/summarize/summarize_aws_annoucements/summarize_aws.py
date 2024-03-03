import xml.etree.ElementTree as ET
from bs4 import BeautifulSoup
import requests
import shutil
import time

from datetime import date
from datetime import timedelta
from datetime import datetime
import os
from langchain import PromptTemplate

from pptx import Presentation
from pptx.util import Inches
from pptx.util import Cm, Pt

from datetime import datetime

import gradio as gr
import json
import boto3
import uuid

from botocore.exceptions import ClientError

#for future use
sample_cost_explorer_service_data = """

		{
		"customerName": "customer", 
		"Services": ["EC2", "Lambda", "DMS"]
		}

"""

SENDER = "AWS Annoucements <test@amazon.com>"
RECIPIENT = "test@amazon.com"


# If necessary, replace us-west-2 with the AWS Region you're using for Amazon SES.
AWS_REGION = "us-west-2"

# The subject line for the email.

INTRO_1 = """

<!DOCTYPE html> <html> <head> <title>HTML Table Generator</title> <style> table { border:1px solid #b3adad; border-collapse:collapse; padding:5px; } table th { border:1px solid #b3adad; padding:5px; background: #ffbd2e; color: #313030; } table td { border:1px solid #b3adad; text-align:left; padding:5px; background: #ffffff; color: #313030; } </style> </head> <body> 

Hello Team!<br>

<p>On behalf of your AWS Enterprise Account Team, I present 🎁 summary of launches 🚀 for requested date. This is a limited and discretionary disclosure to Enterprise Support customers under NDA.<br></p>

<p>Follow the <a href="https://aws.amazon.com/blogs/aws/">AWS Blog</a>, <a href="https://aws.amazon.com/new/">What's New from Amazon Web Services</a>, and <a href="https://aws.amazon.com/podcasts/aws-podcast/">AWS Podcasts</a> for up-to-date news and information.<br></p>

<p>Please contact us if you have any questions or need more information on any of these services.<br></p>

<b>
"""

INTRO_3 = """

<p><a href="https://aws.amazon.com/new/">What's New at AWS - Cloud Innovation & News</a></p>
</b>

"""
            
# The HTML body of the email.        

BODY_HTML_PART1 = """<font size="2px"><table> <thead> <tr> <th>Title</th> <th>Summary (Generated from LLM)</th> <th>Link</th> <th>Date</th> </tr> </thead> <tbody>"""

BODY_HTML_PART2 = """</tbody> </table></font> </body> </html>"""


# The character encoding for the email.
CHARSET = "UTF-8"

# Create a new SES resource and specify a region.
email_client = boto3.client('ses',region_name=AWS_REGION)

bedrock_client = boto3.client('bedrock' , 'us-west-2', endpoint_url = 'https://bedrock.us-west-2.amazonaws.com')

dynamodb = boto3.resource('dynamodb')
table = dynamodb.Table('llm_track')

#function to save prompt and its responses to dynamodb

def savetoDDB(type,prompt, response_llm,parsedDate,link,typeOfExtract):
	now = datetime.now()
	date_time = now.strftime("%m/%d/%Y, %H:%M:%S")

	response = table.put_item(
			Item={
				'llm_type': type,
				'insert_time': date_time,
				'prompt' : prompt,
				'response' : response_llm,
				'parsedDate' : parsedDate,
				'link' : link,
				'typeOfExtract' : typeOfExtract
			}
		)
	
	status_code = response['ResponseMetadata']['HTTPStatusCode']
	return status_code

#use LLM with different types

def interactWithLLM(prompt,type):

	if type == 'titan':
		print("**THE LLM TYPE IS -->" + type)
		#Test for invoke model begins
		parameters = {
			"maxTokenCount":512,
			"stopSequences":[],
			"temperature":0,
			"topP":0.9
		}
		body = json.dumps({"inputText": prompt, "textGenerationConfig": parameters})
		modelId = "amazon.titan-tg1-large" #"amazon.titan-tg1-large"
		accept = "application/json"
		contentType = "application/json"

		response = bedrock_client.invoke_model(
			body=body, modelId=modelId, accept=accept, contentType=contentType
		)

		response_body = json.loads(response.get("body").read())

		response_text_titan = response_body.get("results")[0].get("outputText")

		return response_text_titan
	
	elif type == 'claude':
		print("**THE LLM TYPE IS -->" + type)
		body = json.dumps({"prompt": prompt,
                 "max_tokens_to_sample":300,
                 "temperature":1,
                 "top_k":250,
                 "top_p":0.999,
                 "stop_sequences":[]
                  }) 
		modelId = 'anthropic.claude-v2' # change this to use a different version from the model provider
		accept = 'application/json'
		contentType = 'application/json'
		print("prompt---->" + prompt)
		response = bedrock_client.invoke_model(body=body, modelId=modelId, accept=accept, contentType=contentType)
		response_body = json.loads(response.get('body').read())

		response_text_claude = response_body.get('completion')

		return response_text_claude

#main worker function
def fetchAWSWhatsNew(enteredDate,email,llm_type):

	print("llm_type----" + str(llm_type))
	print("enteredDate----" + str(enteredDate))
	print("email----" + str(email))

	if str(email) == '' : 
		email = RECIPIENT

	if str(llm_type) == 'None' : 
		yield "Please select LLM type"
		return

	#print(f"llm_type as {type(llm_type)} is {llm_type}")
	enteredDateTimeObj = datetime.strptime(enteredDate, "%d %b %Y") #'%Y-%m-%d %H:%M:%S' Thu, 17 Aug 2023 21:32:37 +0000%m-%Y") #'%Y-%m-%d %H:%M:%S' Thu, 17 Aug 2023 21:32:37 +0000
		
	#print(f"Date as {type(enteredDateTimeObj)} is {enteredDateTimeObj}")
	parsedEnteredDate = str(enteredDateTimeObj.day)+"-"+str(enteredDateTimeObj.month)+"-"+str(enteredDateTimeObj.year)	

	prompt_summary_for_email_claud  ="""
Human:  Here are the details of annoucements:
<annoucement>
{text}
</annoucement>
Summarize the above annoucement in one sentence
Assistant: 
	"""

	prompt_summary_for_email_titan  = """	
	{text}
Summarize the above in one sentence
	"""  

	prompt_summary_for_email = ""

	prompt_takeaways = """
	{text}
Give 2 takeaways from the above text
	""" 

	
#	prompt_takeaways_tech_claude = """{text}
#Give 2 takeaways from the above text and list technology names
#	""" 

	prompt_takeaways_tech_claude = """
Human:  Here are the details of annoucements:
<annoucements>
{text}
</annoucements>
Give 2 key takeaways from the above annoucements and list technology names also. In the output takeaways under "Takeaways:" and list technology names under "Technology names:"
Assistant: 
	"""

	prompt_technology_names = """
	{text}

	list technology names in above text
	"""

	url = 'https://aws.amazon.com/about-aws/whats-new/recent/feed/'
	response = requests.get(url, stream=True)
	access_url_message = "Accessing " + url
	for k in range(len(access_url_message)):
		time.sleep(0.01)
		yield access_url_message[: k+1]
	#yield "Accessing " + url
	time.sleep(1)

	with open('aws_whats_new.xml', 'wb') as out_file:
		shutil.copyfileobj(response.raw, out_file)

	prep_start_message = "The latest annoucements file successfully downloaded and now processing with start"
	for k in range(len(prep_start_message)):
		time.sleep(0.01)
		yield prep_start_message[: k+1]
	#yield "The latest annoucements file successfully downloaded and now processing with start"
	time.sleep(1)


	tree = ET.parse('aws_whats_new.xml')
	root = tree.getroot()

	#print(root)

	#print(root.tag)


	#for item in root.iter('item'):
		#print(item)

	prs=Presentation()

	lyt=prs.slide_layouts[0] # choosing a slide layout
	slide=prs.slides.add_slide(lyt) # adding a slide

	title=slide.shapes.title # assigning a title

	subtitle=slide.placeholders[1] # placeholder for subtitle

	title.text="AWS whats new" # title
	subtitle.text="All AWS annoucments for date " + enteredDate # subtitle

	enteredDateCount = 0

	for item in root.findall('channel/item'):
		pubDate =  item.find('pubDate').text
		pubDate =  pubDate.split(',')[1] #<pubDate>Fri, 15 Sep 2023 17:25:53 +0000</pubDate>

		pubDate =  pubDate.split('+')[0].strip() #15 Sep 2023 17:25:53 +0000

		dateTimeObj = datetime.strptime(pubDate, "%d %b %Y %H:%M:%S") #'%Y-%m-%d %H:%M:%S' Thu, 17 Aug 2023 21:32:37 +0000%m-%Y") #'%Y-%m-%d %H:%M:%S' Thu, 17 Aug 2023 21:32:37 +0000
		
		#convert the <pubDate> to the entered data format (DD MON YYYY)
		parsedDate = str(dateTimeObj.day)+" "+str(dateTimeObj.strftime("%b"))+" "+str(dateTimeObj.year)		
		#print("parsedDate from xml for count--"+parsedDate)
		#print("enteredDate for count--"+enteredDate)
		if enteredDate == parsedDate:
			enteredDateCount = enteredDateCount + 1

	print("enteredDateCount--"+str(enteredDateCount))
	if enteredDateCount == 0 :
		print("No entries found")
		no_annoucments_found_message = "No annoucments found for the data " + enteredDate
		for k in range(len(no_annoucments_found_message)):
			time.sleep(0.01)
			yield no_annoucments_found_message[: k+1]
		#yield "No annoucments found for the data " + enteredDate
		return
	prepare_slides_message =  "Found " + str(enteredDateCount) + " annoucements 📣 made on " + enteredDate + " getting ready to prepare the slides..."
	for k in range(len(prepare_slides_message)):
		time.sleep(0.01)
		yield prepare_slides_message[: k+1]
	#yield "Found " + str(enteredDateCount) + " annoucements made on " + enteredDate + " getting ready to prepare the slides..."
	time.sleep(2)

	i=1
	
	processing_start_message = "Processing starts for all the annoucements made on " + enteredDate
	for k in range(len(processing_start_message)):
		time.sleep(0.01)
		yield processing_start_message[: k+1]
	#yield "Processing starts for all the annoucements made on " + enteredDate
	annoucements_table = ""
	for item in root.findall('channel/item'):
		link = item.find('link').text
		title = item.find('title').text
		pubDate =  item.find('pubDate').text
		itemDescription =  item.find('description').text
		

		pubDate =  pubDate.split(',')[1]
		#print("pubDate--"+pubDate)

		pubDate =  pubDate.split('+')[0].strip()

		dateTimeObj = datetime.strptime(pubDate, "%d %b %Y %H:%M:%S") #'%Y-%m-%d %H:%M:%S' Thu, 17 Aug 2023 21:32:37 +0000%m-%Y") #'%Y-%m-%d %H:%M:%S' Thu, 17 Aug 2023 21:32:37 +0000
		
		#print(f"Date as {type(dateTimeObj)} is {dateTimeObj}")
		
		parsedDate = str(dateTimeObj.day)+"-"+str(dateTimeObj.month)+"-"+str(dateTimeObj.year)	
		#print("parsedDate--"+parsedDate)

		if parsedEnteredDate == parsedDate:

			Second_Layout = prs.slide_layouts[5]
			second_slide = prs.slides.add_slide(Second_Layout)

			second_slide.shapes.title.text = title
			title_para = second_slide.shapes.title.text_frame.paragraphs[0]
			title_para.font.size = Pt(18)

			#
			#print(link)
			r  = requests.get(link)

			data = r.text

			soup = BeautifulSoup(data,'html.parser')

			final_text = ''
			nestedDiv_list = soup.find_all('div',{"class": "aws-text-box"})

			for lists in nestedDiv_list:
				element_checked = lists.find('p')
				if element_checked!=None:
					text_got = lists.find('p').text
					final_text = final_text + text_got
					#print(text_got)

			attempt_summarize_message = "[SLIDE # " + str(i) + " of " + str(enteredDateCount) + " SUMMARIZE ATTEMPT 🤖] : LLM will now attempt to summarize the details from new annoucment and extract technology names " + "'" + title + "'" + " annouced on on " + parsedDate + " to start preparing for slide # "+ str(i)
			#yield "LLM will now attempt to summarize the details from new annoucment " + "'" + title + "'" + " annouced on on " + parsedDate + " to start creating slide # "+ str(i)
			for k in range(len(attempt_summarize_message)):
				time.sleep(0.00)
				yield attempt_summarize_message[: k+1]
			
			#print("final_text before strip-->" + final_text)
			final_text = final_text.strip()
			#print("final_text after strip-->" + final_text)
			
			#as of writing this app, claude was able to work off one prompt and titan takes two
			if llm_type == 'claude':

				prompt_template_for_summary_generate = PromptTemplate.from_template(prompt_takeaways_tech_claude)
				prompt_data_for_summary_generate = prompt_template_for_summary_generate.format(text=final_text)

				response_text = interactWithLLM(prompt_data_for_summary_generate,llm_type)
				savetoDDB(llm_type,prompt_data_for_summary_generate,response_text,parsedDate,link,"summarize_and_tech") #savetoDDB(type,prompt, response_llm,cacheKey1,cacheKey2):

			elif llm_type == 'titan':
				
				prompt_template_for_summary_generate = PromptTemplate.from_template(prompt_takeaways)
				prompt_data_for_summary_generate = prompt_template_for_summary_generate.format(text=final_text)

				response_text = interactWithLLM(prompt_data_for_summary_generate,llm_type)
				savetoDDB(llm_type,prompt_data_for_summary_generate,response_text,parsedDate,link,"summarize")
			#response_text = llm.predict(prompt_data_for_summary_generate) #return a response to the prompt

			llm_done_summarizing_message = "[SLIDE # " + str(i) + " of " + str(enteredDateCount) +" SUMMARIZE DONE  ✅] : LLM done summarizing the details for " + "'" + title + "'"
			for k in range(len(llm_done_summarizing_message)):
				time.sleep(0.00)
				yield llm_done_summarizing_message[: k+1]
			#yield "LLM done summarizing the details"
			time.sleep(1)

			attempt_extracting_tech_names_message = "[SLIDE # " + str(i) + " of " + str(enteredDateCount) + " TECH EXTRACT ATTEMPT 🤖] : LLM will now attempt to extract technology names from the details for " + "'" + title + "'"
			for k in range(len(attempt_extracting_tech_names_message)):
				time.sleep(0.00)
				yield attempt_extracting_tech_names_message[: k+1]
			#yield "LLM done summarizing the details"
			time.sleep(1)
			#yield "LLM will now attempt to extract technology names from the details"

			response_text_tech = ''

			if llm_type == 'titan':
				prompt_template_for_tech_generate = PromptTemplate.from_template(prompt_technology_names)
				prompt_data_for_tech_generate = prompt_template_for_tech_generate.format(text=final_text)
				prompt_summary_for_email = prompt_summary_for_email_titan

				response_text_tech = interactWithLLM(prompt_data_for_tech_generate,llm_type)
				savetoDDB(llm_type,prompt_data_for_tech_generate,response_text_tech,parsedDate,link,"tech")
			elif llm_type == 'claude':
				response_text_tech = ''
				prompt_summary_for_email = prompt_summary_for_email_claud

			#response_text_tech = llm.predict(prompt_data_for_tech_generate) #return a response to the prompt
			#print("summary text---" + response_text)
			done_extracting_tech_names_message = "[SLIDE # " + str(i) + " of " + str(enteredDateCount) + " TECH EXTRACT DONE ✅] : LLM done extract technology name for " + "'" + title + "'"
			for k in range(len(done_extracting_tech_names_message)):
				time.sleep(0.00)
				yield done_extracting_tech_names_message[: k+1]
			#yield "LLM done extract technology name"
			time.sleep(1)

			creating_slide_message =  "[SLIDE # " + str(i) + " of "  + str(enteredDateCount) + "  FINALIZING 🛠️ ] : Creating slide # "+ str(i) + " and reading new annoucment " + "'" + title + "'" + " annouced on on " + parsedDate + " to start creating slide # "+ str(i)
			for k in range(len(creating_slide_message)):
				time.sleep(0.00)
				yield creating_slide_message[: k+1]

			time.sleep(1)
			#yield "Creating slide # "+ str(i) + " and reading new annoucment " + "'" + title + "'" + " annouced on on " + parsedDate + " to start creating slide # "+ str(i)

			textbox = second_slide.shapes.add_textbox(Inches(1), Inches(1.5),Inches(9), Inches(5))
			textframe = textbox.text_frame
			textframe.word_wrap = True 
			paragraph = textframe.add_paragraph()
			paragraph.text = response_text + "\n\n" + response_text_tech + "\n\n" + link
			paragraph.font.size = Pt(18)

			print("completed slide #"+ str(i))

			prompt_template_for_summary_generate_email = PromptTemplate.from_template(prompt_summary_for_email)
			prompt_data_for_summary_generate_email = prompt_template_for_summary_generate_email.format(text=final_text)
			email_summary_response = interactWithLLM(prompt_data_for_summary_generate_email,llm_type)
			print("email_summary_response---->" + email_summary_response)
			annoucements_table = annoucements_table + "<tr>"+ "<td>"+title+"</td>" + "<td>"+email_summary_response+"</td><td>"+link+"</td><td>"+pubDate+"</td></tr>"
			
			i=i+1

	print("i----"+str(i))
	ppt_name = "aws_whats_new" + "_" + enteredDate + "_" + llm_type + "_" + str(uuid.uuid1()) + ".pptx"
	prs.save(ppt_name) # saving file
	
	response_text_final = "Congratulations 🎉 , you have successfully created " + str(i) + " slides (1 extra slide for intro) for your AWS announcement. You can find them in your local folder under the name " + ppt_name + ". An email also has been sent which can be used by TAM or account team to forward to your customers "
	for k in range(len(response_text_final)):
		time.sleep(0.01)
		yield response_text_final[: k+1]
	yield response_text_final

	email_body = INTRO_1 + "<p>Launch recap from " + enteredDate +  "("+str(enteredDateCount)+" items):</p>" + INTRO_3 + BODY_HTML_PART1 + annoucements_table + BODY_HTML_PART2
	email_subject = "AWS Annoucements 📣 on " + enteredDate
	print("email_body--" + email_body)

	# Send the email via SES.
	try:
		#Provide the contents of the email.
		response = email_client.send_email(
			Destination={
				'ToAddresses': [
					email,
				],
			},
			Message={
				'Body': {
					'Html': {
						'Charset': CHARSET,
						'Data': email_body,
					},
					'Text': {
						'Charset': CHARSET,
						'Data': email_body,
					},
				},
				'Subject': {
					'Charset': CHARSET,
					'Data': email_subject,
				},
			},
			Source=SENDER,
			# If you are not using a configuration set, comment or delete the
			# following line
			#ConfigurationSetName=CONFIGURATION_SET,
		)
	# Display an error if something goes wrong.	
	except ClientError as e:
		print(e.response['Error']['Message'])
	else:
		print("Email sent! Message ID:"),
		print(response['MessageId'])


#When app loads, fetch latest RSS feed content for the AWS whats new

url = 'https://aws.amazon.com/about-aws/whats-new/recent/feed/'
response = requests.get(url, stream=True)

#save the fetched rss xml to a local file
with open('aws_whats_new.xml', 'wb') as out_file:
	shutil.copyfileobj(response.raw, out_file)

tree = ET.parse('aws_whats_new.xml')
root = tree.getroot()

#Fetch all dates in the RSS feed downloaded xml file to put on the UI as examples

masterDatesArray = []
allDates = []

for item in root.findall('channel/item'):
	link = item.find('link').text
	title = item.find('title').text
	pubDate =  item.find('pubDate').text
		
	pubDate =  pubDate.split(',')[1]

	pubDate =  pubDate.split('+')[0].strip()

	dateTimeObj = datetime.strptime(pubDate, "%d %b %Y %H:%M:%S") #'%Y-%m-%d %H:%M:%S' Thu, 17 Aug 2023 21:32:37 +0000%m-%Y") #'%Y-%m-%d %H:%M:%S' Thu, 17 Aug 2023 21:32:37 +0000
				
	parsedDate = str(dateTimeObj.day)+" "+str(dateTimeObj.strftime("%b"))+" "+str(dateTimeObj.year)	
	#print("parsedDate--"+parsedDate)
	if parsedDate not in masterDatesArray:
		parsedDateSet = []
		parsedDateSet.insert(0,parsedDate)
		allDates.append(parsedDateSet)
	masterDatesArray.append(parsedDate)


#allDates = list(allDates)

#sort all the dates to show from latest date to oldest on the UI in terms of examples

#allDates.sort(key=lambda date: datetime.strptime(date, "%d %b %Y"),reverse=True)
print(allDates)

#also show a pre-populated date which is current date minus one to show a day before releases or in case if this run when hour changes
today = date.today()-timedelta(days = 1)

# Get month, day and year. This is to populate a default date 
d2 = today.strftime("%-d %b %Y")
print("d2 =", d2)

label = "Enter the date :"
lines = 2
info = "Please enter a date to get all the annoucements for that date. This is version1 of demo where only one date is entered in the format (DD MON YYYY) example below. By default yesterday's date is pre-filled"
title = "📣 AWS Announcements Assistant 🤖"
description = "### Evaluate 🔎 AWS Announcements and auto-generate a presentation. This system will fetch 🎣 latest AWS annoucements, using LLM will summarize the details and extract technology names from the details and auto-generate a presentation for you. The system also sends out annoucements email with summary generated from LLM"


#Gradio interface to instantiate the UI controls.
demo = gr.Interface(
    fetchAWSWhatsNew,
	[gr.Textbox(value = d2,label = label,lines=lines,info=info),
  	gr.Textbox(label = "Enter email :", value = "test@test.com",info="Enter your email to receive the summarized annoucement email. Only verified emails will receive communication"),
  	 gr.Radio(["titan", "claude"], label="LLM Types", info="Select from available LLMs")],
	"text",
    title = title,
    description = description,
    examples=allDates
)

#demo.queue().launch()
#demo.queue().launch(share=True)
demo.queue().launch(server_port=8080)




